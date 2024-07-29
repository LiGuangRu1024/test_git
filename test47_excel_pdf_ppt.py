# @time     ：2024/7/15 17:08
# @author   : 莉光哈哈哈
# @file     : test47_excel_pdf_ppt.py
# @software : PyCharm
'''
excel文件处理
'''
from openpyxl import load_workbook
import pandas as pd
from openpyxl.workbook import Workbook

wb = load_workbook(filename='test47_example.xlsx')
sheet = wb.active
cell_value = sheet['A1'].value
print(cell_value)

df = pd.read_excel('test47_example.xlsx')
print(df.head())

wb = Workbook()
ws = wb.active
ws['A1'] = 'Hello'
wb.save('output.xlsx')

data = {'Name': ['john', 'anna'], 'Age': [24, 35]}
df = pd.DataFrame(data)
df.to_excel('output.xlsx', index=False)
# ---------------------------------------------------------------------------------------------------
'''
pdf文件处理
'''
import PyPDF2
from pdfminer.high_level import extract_text

with open('test47_example.pdf', 'rb') as file:
    pdf_reader = PyPDF2.PdfFileReader(file)
    page = pdf_reader.getPage(0)
    text = page.extractText()
    print(text)

text = extract_text("test47_example.pdf")
print(text)
# -----------------------------------------------------------------------------------------------------
'''
ppt文件处理
'''
from pptx import Presentation

# 加载ppt
prs = Presentation('test47_example.pptx')

# 访问第一个幻灯片
slide = prs.slides[0]

# 添加文本框
txBox = slide.shapes.add_textbox(left, top, width, heigh)
tf = txBox.text_frame
tf.text = 'hello world'

# 保存更改
prs.save('output.pptx')
